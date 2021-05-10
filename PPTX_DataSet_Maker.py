from pptx import Presentation
from pptx.util import Inches
import os

'''
    개발자 : 김민규  
    연락처 : rlaalsrb4175@gmail.com
    작업일 : 2020.12.30
    설명 : 대용량 pptx 이미지 작업을 자동화하는 코드
    사용하기 앞서 데이터 분류 작업을 잘 해놓을 것
    D:\김민규\최은정교수님\PPT_H
    위 폴더를 참고하기 바람
    출력되는 ppt의 내용은 '환자번호 - coronal - 사진들 - transverse - 사진들'로 반복됨
    실행 시키려면 F5를 눌러서 상단에 나오는 리스트 중에 Python File클릭
    이후 수정사항 반영은 ctrl + s 로 저장
'''


class make_ppt() :
    target = ""
    def pre_view(self,target_dir):
        global target
        target = target_dir
        dir_list = os.listdir(target_dir)
        print(target)
        print(dir_list)
        return dir_list


    def ppt_work(self, save_dir):
        prs = Presentation()
        dir_list = os.listdir(target)  # 읽어오고자 하는 폴더위치 삽입
        print(dir_list)

        blank_slide_layout = prs.slide_layouts[6]  # 이미지를 넣을 슬라이드 생성
        title_slide_layout = prs.slide_layouts[0]  # 텍스트를 넣을 슬라이드 생성
        text_slide = prs.slides.add_slide(title_slide_layout)
        text_slide.shapes.title.text = target.split("\\")[-1]
        for slidetitle in dir_list:  # case의 숫자만큼 반복
            dir_path = target + '/' + slidetitle
            img_list = os.listdir(dir_path)  # 이미지 파일 불러오기
            text_slide = prs.slides.add_slide(title_slide_layout)
            text_slide.shapes.title.text = slidetitle  # case 번호를 보여주는 슬라이드
            for image in img_list:  # 이미지 설정

                left = Inches(1)
                width = Inches(7.5)
                height = Inches(7.5)
                top = Inches(0)

                image_slide = prs.slides.add_slide(blank_slide_layout)
                img_path = dir_path + '/' + image
                pic = image_slide.shapes.add_picture(img_path, left, top, width=width, height=height)
        direction = save_dir + '/' + 'DataSet.pptx'
        prs.save(direction)  # 파일 이름 (확장자 반드시 포함), D:\김민규\Python 이 경로로 저장됨