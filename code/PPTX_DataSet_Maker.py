from pptx import Presentation
from pptx.util import Inches
import os
from tqdm import tqdm
import tkinter as tk
from tkinter import messagebox as msg

class make_ppt() :
    target = ""
    def pre_view(self,target_dir):
        global target
        target = target_dir
        dir_list = os.listdir(target_dir)
        print(target)
        print(dir_list)
        return dir_list


    def ppt_work(self, save_dir):
        prs = Presentation()
        dir_list = os.listdir(target)  # 읽어오고자 하는 폴더위치 삽입
        print(dir_list)

        blank_slide_layout = prs.slide_layouts[6]  # 이미지를 넣을 슬라이드 생성
        title_slide_layout = prs.slide_layouts[0]  # 텍스트를 넣을 슬라이드 생성
        text_slide = prs.slides.add_slide(title_slide_layout)
        text_slide.shapes.title.text = target.split("\\")[-1]

        for slidetitle in dir_list:  # case의 숫자만큼 반복
            dir_path = target + '/' + slidetitle
            img_list = os.listdir(dir_path)  # 이미지 파일 불러오기
            text_slide = prs.slides.add_slide(title_slide_layout)
            text_slide.shapes.title.text = slidetitle  # case 번호를 보여주는 슬라이드

            for image in tqdm(img_list):  # 이미지 설정

                left = Inches(1)
                width = Inches(7.5)
                height = Inches(7.5)
                top = Inches(0)

                image_slide = prs.slides.add_slide(blank_slide_layout)
                img_path = dir_path + '/' + image

        direction = save_dir + '/' + 'DataSet.pptx'
        prs.save(direction)  # 파일 이름 (확장자 반드시 포함)

        # root = Tk()
        # root.withdraw()
        msg.showinfo('Message', '생성 완료')